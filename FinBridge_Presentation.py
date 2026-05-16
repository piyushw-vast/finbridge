from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Colour palette ──────────────────────────────────────────────────────────
INDIGO   = RGBColor(0x4f, 0x46, 0xe5)
INDIGO_D = RGBColor(0x31, 0x2e, 0x81)
WHITE    = RGBColor(0xff, 0xff, 0xff)
SLATE    = RGBColor(0x1e, 0x29, 0x3b)
SLATE_L  = RGBColor(0x64, 0x74, 0x8b)
EMERALD  = RGBColor(0x10, 0xb9, 0x81)
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)
ROSE     = RGBColor(0xf4, 0x3f, 0x5e)
LIGHT_BG = RGBColor(0xf8, 0xf9, 0xff)
CARD_BG  = RGBColor(0xff, 0xff, 0xff)
BORDER   = RGBColor(0xe2, 0xe8, 0xf0)

BLANK = prs.slide_layouts[6]  # completely blank

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=12, bold=False, color=SLATE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bg(slide, color=LIGHT_BG):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def indigo_bar(slide, h=0.08):
    add_rect(slide, 0, 7.5 - h, 13.33, h, fill=INDIGO)

def slide_label(slide, text):
    add_text(slide, text, 0.45, 0.22, 4, 0.3, size=8, bold=True,
             color=INDIGO, align=PP_ALIGN.LEFT)

def card(slide, l, t, w, h, title=None, title_color=SLATE, body_lines=None,
         accent=None, icon_char=None):
    add_rect(slide, l, t, w, h, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    if accent:
        add_rect(slide, l, t, 0.04, h, fill=accent)
    if title:
        add_text(slide, title, l + 0.18, t + 0.12, w - 0.28, 0.28,
                 size=10, bold=True, color=title_color)
    if body_lines:
        y = t + 0.42 if title else t + 0.15
        for line in body_lines:
            add_text(slide, line, l + 0.18, y, w - 0.28, 0.24, size=8.5, color=SLATE_L)
            y += 0.24

def pill(slide, text, l, t, w=1.6, color=INDIGO, text_color=WHITE, h=0.28):
    add_rect(slide, l, t, w, h, fill=color)
    add_text(slide, text, l, t + 0.02, w, h - 0.04, size=8, bold=True,
             color=text_color, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=INDIGO_D)
add_rect(sl, 0, 0, 13.33, 7.5, fill=INDIGO)
# diagonal accent
add_rect(sl, 9.5, 0, 4, 7.5, fill=RGBColor(0x43, 0x38, 0xca))
add_rect(sl, 10.8, 0, 3, 7.5, fill=RGBColor(0x38, 0x2d, 0xb5))

# logo box
add_rect(sl, 1.0, 2.0, 0.7, 0.7, fill=WHITE)
add_text(sl, "FB", 1.0, 2.0, 0.7, 0.7, size=18, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

add_text(sl, "FinBridge", 1.85, 1.95, 5, 0.6, size=32, bold=True, color=WHITE)
add_text(sl, "A Multi-Tenant Financial Data Exchange Platform",
         1.85, 2.62, 7, 0.4, size=14, color=RGBColor(0xc7, 0xd2, 0xfe))
add_rect(sl, 1.0, 3.18, 4, 0.04, fill=RGBColor(0x81, 0x8c, 0xf8))
add_text(sl, "AI-powered invoice scanning  ·  Multi-role workflows  ·  Real-time intelligence",
         1.0, 3.35, 8, 0.35, size=11, color=RGBColor(0xa5, 0xb4, 0xfc), italic=True)

add_text(sl, "Team: Token Heroes", 1.0, 5.6, 5, 0.35, size=12, bold=True, color=WHITE)
add_text(sl, "VAST Hackathon 2026  ·  May 16", 1.0, 6.0, 5, 0.3, size=10,
         color=RGBColor(0xa5, 0xb4, 0xfc))
add_text(sl, "github.com/piyushw-vast/finbridge", 1.0, 6.35, 6, 0.25, size=9,
         color=RGBColor(0x81, 0x8c, 0xf8))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "01  /  PROBLEM")

add_text(sl, "How Businesses Share Financial Data Today",
         0.5, 0.5, 10, 0.55, size=24, bold=True, color=SLATE)
add_text(sl, "Mid-sized companies exchange invoices, payments and salary registers with accounting firms through:",
         0.5, 1.18, 12, 0.35, size=11, color=SLATE_L)

chaos = [
    ("📧", "Email", "Attachments lost in threads, no version control"),
    ("💬", "WhatsApp", "Informal, unsecured, no audit trail"),
    ("📁", "Shared Drives", "Manual naming, zero extraction, no status tracking"),
]
x = 0.5
for icon, title, desc in chaos:
    add_rect(sl, x, 1.75, 3.8, 1.4, fill=RGBColor(0xff, 0xf1, 0xf2), line=ROSE, line_w=Pt(0.75))
    add_rect(sl, x, 1.75, 3.8, 0.04, fill=ROSE)
    add_text(sl, icon + "  " + title, x + 0.15, 1.88, 3.5, 0.35, size=13, bold=True, color=ROSE)
    add_text(sl, desc, x + 0.15, 2.28, 3.5, 0.7, size=9.5, color=SLATE_L)
    x += 4.1

add_text(sl, "The result:", 0.5, 3.42, 3, 0.3, size=11, bold=True, color=SLATE)
pains = [
    "Accountants manually re-key data into Zoho / Tally / QuickBooks",
    "Errors, duplicates and fraud go undetected",
    "No visibility for companies on invoice status",
    "Compliance and GST reconciliation done by hand",
]
y = 3.42
for p in pains:
    add_rect(sl, 3.5, y + 0.06, 0.18, 0.18, fill=ROSE)
    add_text(sl, p, 3.85, y, 8.8, 0.3, size=10.5, color=SLATE)
    y += 0.38

add_rect(sl, 0.5, 5.3, 12.3, 0.85, fill=RGBColor(0xef, 0xf6, 0xff), line=INDIGO, line_w=Pt(1))
add_text(sl, "💡  Opportunity:  Replace ad-hoc sharing with a structured, AI-powered platform that extracts, validates and routes financial documents automatically.",
         0.75, 5.38, 11.8, 0.65, size=10.5, color=INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "02  /  SOLUTION")

add_text(sl, "FinBridge — What We Built",
         0.5, 0.5, 10, 0.55, size=24, bold=True, color=SLATE)

# 3 tenant levels
tenants = [
    (INDIGO, "Platform Admin", "Onboards accounting firms\nand their admins"),
    (RGBColor(0x06, 0x95, 0x6a), "Accounting Firm Admin", "Manages accountants\nOnboards client companies\nConfigures payment heads"),
    (RGBColor(0x0e, 0x7a, 0xb5), "Company Admin / Users", "Uploads invoices, payments\nbank statements, salary registers"),
    (RGBColor(0x92, 0x40, 0x0e), "Accountants", "Reviews, corrects and accepts\ntransactions. Publishes MIS reports"),
]
x = 0.4
for color, title, desc in tenants:
    add_rect(sl, x, 1.3, 2.9, 1.5, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    add_rect(sl, x, 1.3, 2.9, 0.06, fill=color)
    add_text(sl, title, x + 0.14, 1.42, 2.65, 0.32, size=10.5, bold=True, color=color)
    add_text(sl, desc, x + 0.14, 1.8, 2.65, 0.85, size=9, color=SLATE_L)
    if x < 9:
        add_text(sl, "→", x + 2.92, 1.9, 0.3, 0.35, size=14, bold=True, color=SLATE_L, align=PP_ALIGN.CENTER)
    x += 3.1

# Core features
add_text(sl, "Core Capabilities", 0.5, 3.05, 6, 0.32, size=13, bold=True, color=SLATE)
feats = [
    (INDIGO,   "AI Bill Scanning",          "Upload any invoice image or PDF. Groq Vision + PyMuPDF + OCR extract all fields in seconds with confidence scoring"),
    (EMERALD,  "Accountant Review Workflow", "Review queue, field-level corrections, accept/reject with notes — full audit trail on every action"),
    (AMBER,    "MIS Reports",               "Accountants publish reports. Companies view inline summaries and download as PDF. GST split included"),
    (ROSE,     "GST Dashboard",             "Monthly CGST/SGST/IGST trend, vendor analytics, payment aging buckets, fraud signal detection"),
]
x = 0.4
for color, title, desc in feats:
    add_rect(sl, x, 3.5, 2.9, 1.65, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    add_rect(sl, x, 3.5, 0.05, 1.65, fill=color)
    add_text(sl, title, x + 0.2, 3.6, 2.6, 0.3, size=9.5, bold=True, color=color)
    add_text(sl, desc, x + 0.2, 3.96, 2.6, 1.1, size=8.5, color=SLATE_L)
    x += 3.1

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "03  /  AI EXTRACTION ENGINE")

add_text(sl, "The AI Magic — How Invoice Scanning Works",
         0.5, 0.5, 11, 0.55, size=24, bold=True, color=SLATE)

steps = [
    (INDIGO,   "1",  "Upload",            "Company uploads invoice\nimage or PDF (any quality)"),
    (RGBColor(0x0e, 0x7a, 0xb5), "2", "Vision Enhancement", "Deskew · Denoise · Contrast\nPDF render via PyMuPDF"),
    (RGBColor(0x7c, 0x3a, 0xed), "3", "Multi-Engine Extract", "Groq Vision (LLaMA)\n+ PyMuPDF + Tesseract OCR"),
    (AMBER,    "4",  "Consensus Engine",  "3 engines vote on each field\nConflicts flagged automatically"),
    (EMERALD,  "5",  "Trust Score",       "0–100 risk score assigned\nAuto-accept if score ≥ 85"),
    (ROSE,     "6",  "Review / Accept",   "Accountant reviews conflicts\nCorrects and accepts"),
]

x = 0.35
for color, num, title, desc in steps:
    add_rect(sl, x, 1.3, 2.0, 1.9, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    add_rect(sl, x, 1.3, 2.0, 0.05, fill=color)
    add_rect(sl, x + 0.15, 1.42, 0.38, 0.38, fill=color)
    add_text(sl, num, x + 0.15, 1.4, 0.38, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, x + 0.65, 1.42, 1.25, 0.36, size=9.5, bold=True, color=color)
    add_text(sl, desc, x + 0.15, 1.88, 1.75, 0.85, size=8.5, color=SLATE_L)
    if x < 10.5:
        add_text(sl, "→", x + 2.0, 2.05, 0.32, 0.35, size=14, color=SLATE_L, align=PP_ALIGN.CENTER)
    x += 2.16

# Fields extracted
add_rect(sl, 0.35, 3.45, 12.6, 0.32, fill=RGBColor(0xee, 0xf2, 0xff), line=INDIGO, line_w=Pt(0.5))
add_text(sl, "Fields Extracted:   Vendor Name  ·  GST Number  ·  Invoice Number  ·  Date  ·  Line Items  ·  Subtotal  ·  Tax Amount  ·  Total Amount  ·  Payment Terms",
         0.55, 3.49, 12.2, 0.25, size=8.5, color=INDIGO)

# Trust score breakdown
boxes = [
    (EMERALD, "≥ 85 / 100", "Auto-Accepted", "High confidence — no human review needed"),
    (AMBER,   "60 – 84",    "Human Review",  "Routed to accountant for verification"),
    (ROSE,    "< 60",       "High Risk",     "Mandatory manual review + conflict resolution"),
]
x = 0.35
for color, score, label, desc in boxes:
    add_rect(sl, x, 4.0, 3.9, 1.3, fill=CARD_BG, line=color, line_w=Pt(1.5))
    add_text(sl, score, x + 0.15, 4.1, 1.8, 0.38, size=18, bold=True, color=color)
    add_text(sl, label, x + 0.15, 4.52, 2.5, 0.28, size=10, bold=True, color=SLATE)
    add_text(sl, desc, x + 0.15, 4.82, 3.6, 0.38, size=8.5, color=SLATE_L)
    x += 4.15

add_text(sl, "Duplicate Detection  ·  GST Format Validation  ·  Amount Anomaly Detection  ·  Vendor Trust Scoring",
         0.35, 5.52, 12, 0.28, size=9, color=SLATE_L, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "04  /  ARCHITECTURE")

add_text(sl, "System Architecture",
         0.5, 0.5, 10, 0.55, size=24, bold=True, color=SLATE)

# Frontend box
add_rect(sl, 0.4, 1.2, 5.8, 2.1, fill=RGBColor(0xee, 0xf2, 0xff), line=INDIGO, line_w=Pt(1))
add_text(sl, "React Frontend  (Vite + Tailwind CSS)", 0.65, 1.28, 5.3, 0.3, size=10, bold=True, color=INDIGO)
roles = ["Company Admin / User", "Accountant", "Firm Admin", "Platform Admin"]
rx = 0.55
for r in roles:
    add_rect(sl, rx, 1.72, 1.28, 0.38, fill=INDIGO)
    add_text(sl, r, rx, 1.74, 1.28, 0.34, size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rx += 1.38

# Arrow down
add_text(sl, "REST API  /  JWT Auth", 3.0, 3.45, 3.2, 0.28, size=8.5, color=SLATE_L, align=PP_ALIGN.CENTER)
add_text(sl, "↕", 3.9, 3.2, 0.8, 0.4, size=18, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

# Backend box
add_rect(sl, 0.4, 3.8, 5.8, 2.0, fill=RGBColor(0xf0, 0xfd, 0xf4), line=EMERALD, line_w=Pt(1))
add_text(sl, "FastAPI Backend  (Python + SQLAlchemy async)", 0.65, 3.88, 5.3, 0.3, size=10, bold=True, color=EMERALD)
services = ["Auth + RBAC", "Invoice Engine", "Reports", "Notifications", "Audit Trail"]
rx = 0.55
for s in services:
    add_rect(sl, rx, 4.32, 1.05, 0.35, fill=EMERALD)
    add_text(sl, s, rx, 4.33, 1.05, 0.3, size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rx += 1.12

# DB + AI
add_rect(sl, 0.4, 6.0, 2.7, 0.9, fill=RGBColor(0xff, 0xff, 0xff), line=BORDER, line_w=Pt(0.75))
add_text(sl, "🐘  PostgreSQL", 0.65, 6.08, 2.2, 0.28, size=10, bold=True, color=SLATE)
add_text(sl, "Async · JSONB · Alembic migrations", 0.65, 6.38, 2.2, 0.4, size=7.5, color=SLATE_L)

add_rect(sl, 3.5, 6.0, 2.7, 0.9, fill=RGBColor(0xff, 0xff, 0xff), line=BORDER, line_w=Pt(0.75))
add_text(sl, "⚡  Groq Vision API", 3.75, 6.08, 2.2, 0.28, size=10, bold=True, color=SLATE)
add_text(sl, "LLaMA 3.2 Vision · PyMuPDF · OCR", 3.75, 6.38, 2.2, 0.4, size=7.5, color=SLATE_L)

# Right side — tech stack
add_rect(sl, 7.0, 1.2, 5.9, 5.7, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
add_text(sl, "Tech Stack", 7.25, 1.3, 5.4, 0.32, size=12, bold=True, color=SLATE)

stack = [
    (INDIGO,   "Frontend",   "React 18  ·  Vite  ·  Tailwind CSS  ·  React Router"),
    (EMERALD,  "Backend",    "FastAPI  ·  Python 3.10  ·  SQLAlchemy async  ·  Alembic"),
    (RGBColor(0x0e, 0x7a, 0xb5), "Database", "PostgreSQL 14  ·  asyncpg  ·  JSONB columns"),
    (AMBER,    "AI / ML",    "Groq Vision API  ·  LLaMA 3.2 Vision  ·  PyMuPDF  ·  Tesseract OCR"),
    (ROSE,     "Auth",       "JWT (RS256)  ·  Role-based access control  ·  5 roles"),
    (RGBColor(0x7c, 0x3a, 0xed), "PWA",  "Service Worker  ·  Web App Manifest  ·  Installable on iOS + Android"),
    (SLATE_L,  "DevOps",     "Git  ·  Uvicorn  ·  python-dotenv  ·  Pillow"),
]
y = 1.78
for color, label, tech in stack:
    add_rect(sl, 7.15, y, 0.06, 0.28, fill=color)
    add_text(sl, label, 7.3, y, 1.3, 0.28, size=8.5, bold=True, color=color)
    add_text(sl, tech, 8.65, y, 3.8, 0.28, size=8, color=SLATE_L)
    y += 0.64

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FEATURES CHECKLIST
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "05  /  FEATURES")

add_text(sl, "Everything the Hackathon Asked For — Delivered",
         0.5, 0.5, 12, 0.55, size=24, bold=True, color=SLATE)

# Must have
add_text(sl, "✅  Must-Have Requirements", 0.5, 1.2, 6, 0.32, size=12, bold=True, color=EMERALD)
must = [
    "Multi-tenant role-based auth — Platform Admin → Firm Admin → Company → Accountant",
    "Company onboarding with configurable payment heads & sub-heads by business type",
    "Financial data upload — invoices (purchase/sales), payments, salary registers, bank statements, ledgers",
    "AI-powered bill scanning — Groq Vision + consensus extraction engine + trust scoring",
    "Accountant review workflow — refine, correct, accept, reject with notes",
    "MIS Reports — accountant publishes, company views inline summary + downloads as PDF",
]
y = 1.62
for item in must:
    add_rect(sl, 0.5, y + 0.07, 0.18, 0.18, fill=EMERALD)
    add_text(sl, item, 0.82, y, 5.8, 0.3, size=9, color=SLATE)
    y += 0.36

# Stretch goals
add_text(sl, "🚀  Stretch Goals (Bonus)", 7.0, 1.2, 5.8, 0.32, size=12, bold=True, color=INDIGO)
stretch = [
    "PWA — installable on Android & iOS home screen",
    "Dashboard with insights — cash flow, top vendors, spend by category",
    "Audit trail — every action logged with actor + timestamp",
    "Notifications — real-time bell, per-event types (accept/reject/report)",
    "GST Dashboard — CGST/SGST/IGST split, monthly trend, payment aging",
    "Bulk document upload — bank statements, salary registers (up to 10 files)",
]
y = 1.62
for item in stretch:
    add_rect(sl, 7.0, y + 0.07, 0.18, 0.18, fill=INDIGO)
    add_text(sl, item, 7.32, y, 5.6, 0.3, size=9, color=SLATE)
    y += 0.36

add_rect(sl, 0.4, 4.0, 12.5, 0.04, fill=BORDER)

# Bonus features
add_text(sl, "⭐  Beyond the Brief", 0.5, 4.18, 12, 0.3, size=12, bold=True, color=AMBER)
bonus = [
    ("Risk Radar", "Multi-dimensional risk profile — 6-axis radar chart per invoice"),
    ("Fraud Signals", "Anomaly detection — duplicate invoices, suspicious amounts, weekend dates"),
    ("Vendor Intelligence", "Per-vendor trust score, invoice history, GST validation"),
    ("AI Chat", "Natural language query over your invoice dataset"),
    ("Consensus Heatmap", "Visual disagreement map across 3 extraction engines"),
    ("Dark Mode", "Full dark theme across all pages and components"),
]
x = 0.4
for title, desc in bonus:
    add_rect(sl, x, 4.58, 3.9, 1.3, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    add_rect(sl, x, 4.58, 3.9, 0.05, fill=AMBER)
    add_text(sl, title, x + 0.14, 4.68, 3.6, 0.28, size=9.5, bold=True, color=AMBER)
    add_text(sl, desc, x + 0.14, 5.0, 3.6, 0.7, size=8.5, color=SLATE_L)
    x += 4.15
    if x > 8.5:
        break

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — DEMO FLOW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "06  /  DEMO")

add_text(sl, "End-to-End Demo Flow",
         0.5, 0.5, 10, 0.55, size=24, bold=True, color=SLATE)

flows = [
    (INDIGO,   "Step 1",  "Login as Company Admin\n(companyadmin@acme.com)",   "Dashboard shows 13 invoices,\nspend insights, fraud signals,\npayment aging"),
    (EMERALD,  "Step 2",  "Upload an Invoice\n(any PDF or photo)",              "AI extracts all fields in <5s\nTrust score assigned\nConflicts flagged in real-time"),
    (AMBER,    "Step 3",  "Login as Accountant\n(accountant@demo.com)",         "Review queue shows pending\nCorrect fields, add notes\nAccept or reject"),
    (ROSE,     "Step 4",  "Back to Company\nView Invoice Detail",               "Risk Radar · Invoice Journey\nAI summary · Audit trail\nAll updated live"),
    (RGBColor(0x7c, 0x3a, 0xed), "Step 5", "View MIS Reports\n& GST Dashboard", "Download PDF report\nGST split · monthly trend\nVendor analytics"),
]
x = 0.35
for color, step, action, result in flows:
    add_rect(sl, x, 1.3, 2.38, 2.3, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    add_rect(sl, x, 1.3, 2.38, 0.06, fill=color)
    pill(sl, step, x + 0.5, 1.42, 1.38, color=color)
    add_text(sl, action, x + 0.14, 1.82, 2.1, 0.55, size=9, bold=True, color=SLATE)
    add_rect(sl, x + 0.1, 2.45, 2.18, 0.04, fill=BORDER)
    add_text(sl, result, x + 0.14, 2.55, 2.1, 0.85, size=8.5, color=SLATE_L)
    if x < 9.5:
        add_text(sl, "→", x + 2.38, 2.1, 0.25, 0.35, size=12, bold=True, color=SLATE_L, align=PP_ALIGN.CENTER)
    x += 2.6

# Demo data
add_rect(sl, 0.35, 3.85, 12.6, 0.85, fill=RGBColor(0xee, 0xf2, 0xff), line=INDIGO, line_w=Pt(0.75))
add_text(sl, "Demo Seed Data Ready", 0.6, 3.93, 3, 0.28, size=10, bold=True, color=INDIGO)
add_text(sl, "20 documents  ·  2 companies  ·  5 document types  ·  3 MIS reports  ·  Mix of accepted / under review / rejected / high-risk invoices",
         0.6, 4.22, 12, 0.35, size=9, color=SLATE_L)

add_text(sl, "Login Credentials:", 0.35, 4.9, 3, 0.28, size=10, bold=True, color=SLATE)
creds = [
    ("Company Admin", "companyadmin@acme.com / company123"),
    ("Accountant",    "accountant@demo.com / acc123"),
    ("Firm Admin",    "firmadmin@demo.com / firm123"),
    ("Platform Admin","admin@finbridge.com / admin123"),
]
x = 0.35
for role, cred in creds:
    add_rect(sl, x, 5.25, 3.0, 0.65, fill=CARD_BG, line=BORDER, line_w=Pt(0.5))
    add_text(sl, role, x + 0.12, 5.32, 2.76, 0.24, size=8.5, bold=True, color=SLATE)
    add_text(sl, cred, x + 0.12, 5.56, 2.76, 0.24, size=7.5, color=SLATE_L)
    x += 3.18

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHAT'S NEXT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
indigo_bar(sl)
slide_label(sl, "07  /  ROADMAP")

add_text(sl, "What We'd Build Next",
         0.5, 0.5, 10, 0.55, size=24, bold=True, color=SLATE)

roadmap = [
    (INDIGO,   "Integrations",       "Connect to Zoho Books, Tally, QuickBooks via official APIs — eliminating manual re-entry entirely"),
    (EMERALD,  "Mobile Native App",  "React Native app for field invoice capture — point camera, AI extracts instantly"),
    (AMBER,    "Bank Reconciliation","Automated matching of bank statement entries with uploaded invoices — ML-powered categorization"),
    (ROSE,     "Multi-Currency",     "Forex handling, currency conversion, international vendor support"),
    (RGBColor(0x7c, 0x3a, 0xed), "Scheduled Reports", "Auto-generate and email MIS reports on configurable schedules"),
    (RGBColor(0x0e, 0x7a, 0xb5), "Compliance Engine", "Automated GST filing data prep, TDS tracking, regulatory deadline alerts"),
]
x = 0.4
y_start = 1.25
col = 0
for color, title, desc in roadmap:
    row = col // 2
    c = col % 2
    lx = 0.4 + c * 6.4
    ly = y_start + row * 1.65
    add_rect(sl, lx, ly, 6.1, 1.45, fill=CARD_BG, line=BORDER, line_w=Pt(0.75))
    add_rect(sl, lx, ly, 0.06, 1.45, fill=color)
    add_text(sl, title, lx + 0.22, ly + 0.14, 5.7, 0.3, size=11, bold=True, color=color)
    add_text(sl, desc, lx + 0.22, ly + 0.52, 5.7, 0.72, size=9, color=SLATE_L)
    col += 1

add_rect(sl, 0.4, 6.28, 12.5, 0.65, fill=INDIGO)
add_text(sl, "FinBridge is the foundation. The platform is designed to scale to full accounting automation — replacing email and WhatsApp forever.",
         0.7, 6.36, 12.0, 0.5, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/media/piyushw/C058547C58547360/finbridge/FinBridge_TokenHeroes.pptx"
prs.save(out)
print(f"Saved: {out}")
